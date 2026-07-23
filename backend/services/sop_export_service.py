import base64
from io import BytesIO

from bs4 import BeautifulSoup
from docx import Document
from docx.shared import Inches as DocxInches
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt
from sqlalchemy.orm import Session

from models.sop import SOP

_COVER_FIELD_LABELS = [
    ("project_name", "프로젝트명"),
    ("part_name", "부품명"),
    ("customer", "고객사"),
    ("sample_quantity", "시료 수량"),
    ("remarks", "비고"),
]


def _decode_img(img_tag) -> bytes | None:
    src = img_tag.get("src", "")
    if not src.startswith("data:image"):
        return None
    try:
        _, b64data = src.split(",", 1)
        return base64.b64decode(b64data)
    except Exception:
        return None


def parse_content_html(html: str | None) -> list[dict]:
    """절차 및 방법(content)의 HTML을 문단/목록/이미지 블록 리스트로 변환한다.

    구버전 데이터(순수 텍스트)와 신버전(contentEditable HTML) 모두 지원한다.
    """
    if not html or not html.strip():
        return []

    soup = BeautifulSoup(html, "html.parser")
    if not soup.find(True):
        return [{"type": "paragraph", "text": line} for line in html.split("\n") if line.strip()]

    blocks: list[dict] = []

    def walk(node, ordered: bool | None = None):
        for child in node.children:
            name = getattr(child, "name", None)
            if name == "img":
                data = _decode_img(child)
                if data:
                    blocks.append({"type": "image", "data": data})
            elif name in ("ol", "ul"):
                walk(child, ordered=(name == "ol"))
            elif name == "li":
                text = child.get_text(" ", strip=True)
                if text:
                    blocks.append({"type": "list_item", "text": text, "ordered": bool(ordered)})
                for img_tag in child.find_all("img"):
                    data = _decode_img(img_tag)
                    if data:
                        blocks.append({"type": "image", "data": data})
            elif name in ("p", "div", "span"):
                direct_imgs = child.find_all("img", recursive=False)
                text = child.get_text(" ", strip=True)
                if text:
                    blocks.append({"type": "paragraph", "text": text})
                for img_tag in direct_imgs:
                    data = _decode_img(img_tag)
                    if data:
                        blocks.append({"type": "image", "data": data})
                walk(child, ordered=ordered)
            elif name == "br":
                continue
            elif name is None:
                text = str(child).strip()
                if text:
                    blocks.append({"type": "paragraph", "text": text})
            else:
                walk(child, ordered=ordered)

    walk(soup)
    return blocks


def _split_procedure(content_html: str | None) -> tuple[str, list[bytes]]:
    blocks = parse_content_html(content_html)
    lines: list[str] = []
    images: list[bytes] = []
    counter = 0
    for b in blocks:
        if b["type"] == "image":
            images.append(b["data"])
        elif b["type"] == "list_item":
            if b.get("ordered"):
                counter += 1
                lines.append(f"{counter}. {b['text']}")
            else:
                lines.append(f"- {b['text']}")
        else:
            counter = 0
            lines.append(b["text"])
    return "\n".join(lines), images


def _get_sops_ordered(db: Session, sop_ids: list[int]) -> list[SOP]:
    sops = db.query(SOP).filter(SOP.id.in_(sop_ids)).all()
    order = {sid: i for i, sid in enumerate(sop_ids)}
    sops.sort(key=lambda s: order.get(s.id, 0))
    return sops


# ── DOCX ──────────────────────────────────────────────────
def _add_cover_docx(doc: Document, cover_info: dict | None):
    table = doc.add_table(rows=0, cols=2)
    table.style = "Table Grid"
    for key, label in _COVER_FIELD_LABELS:
        val = (cover_info or {}).get(key)
        if val:
            row = table.add_row()
            row.cells[0].text = label
            row.cells[1].text = val
    doc.add_page_break()


def _add_sop_section_docx(doc: Document, sop: SOP):
    doc.add_heading(f"{sop.sop_number}  {sop.title}", level=2)

    meta_table = doc.add_table(rows=0, cols=2)
    meta_table.style = "Table Grid"
    for label, val in [
        ("버전", sop.version),
        ("분류", sop.category),
        ("최근 개정일", str(sop.revision_date or sop.issue_date or "") or None),
    ]:
        if val:
            row = meta_table.add_row()
            row.cells[0].text = label
            row.cells[1].text = str(val)

    def add_field(label: str, value: str | None):
        if not value:
            return
        doc.add_paragraph(f"■ {label}", style="Heading 3")
        doc.add_paragraph(value)

    add_field("목적", sop.description)
    add_field("시료 수량", sop.sample_quantity)
    add_field("시험 조건", sop.test_condition)

    if sop.test_device:
        doc.add_paragraph("■ 필요 장비", style="Heading 3")
        for line in sop.test_device.split("\n"):
            if line.strip():
                doc.add_paragraph(line.strip(), style="List Bullet")

    if sop.content:
        doc.add_paragraph("■ 절차 및 방법", style="Heading 3")
        for block in parse_content_html(sop.content):
            if block["type"] == "image":
                try:
                    doc.add_picture(BytesIO(block["data"]), width=DocxInches(4))
                except Exception:
                    pass
            elif block["type"] == "list_item":
                style = "List Number" if block.get("ordered") else "List Bullet"
                doc.add_paragraph(block["text"], style=style)
            else:
                doc.add_paragraph(block["text"])

    add_field("판정 기준", sop.judgment_criteria)
    add_field("비고", sop.notes)
    doc.add_page_break()


def generate_docx(db: Session, sop_ids: list[int], variant: str = "절차서", cover_info: dict | None = None) -> bytes:
    sops = _get_sops_ordered(db, sop_ids)
    doc = Document()
    doc.add_heading(f"ES 시험 {variant}", level=1)
    if variant == "계획서":
        _add_cover_docx(doc, cover_info)
    for sop in sops:
        _add_sop_section_docx(doc, sop)
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PPTX ──────────────────────────────────────────────────
def _add_cover_slide_pptx(prs: Presentation, cover_info: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "ES 시험 계획서"
    rows = [(label, (cover_info or {}).get(key, "")) for key, label in _COVER_FIELD_LABELS if (cover_info or {}).get(key)]
    if rows:
        table_shape = slide.shapes.add_table(len(rows), 2, PptxInches(0.5), PptxInches(1.6), PptxInches(9), PptxInches(0.5) * len(rows))
        table = table_shape.table
        for i, (label, val) in enumerate(rows):
            table.cell(i, 0).text = label
            table.cell(i, 1).text = val


def _add_sop_slide_pptx(prs: Presentation, sop: SOP):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(PptxInches(0.3), PptxInches(0.15), PptxInches(9.4), PptxInches(0.5))
    tf = title_box.text_frame
    tf.text = f"{sop.sop_number}  {sop.title}"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True

    fields = [
        ("Scope", sop.description or "-"),
        ("Quantity of Sample", sop.sample_quantity or "-"),
        ("Condition", sop.test_condition or "-"),
        ("Device", sop.test_device or "-"),
        ("Required", sop.judgment_criteria or "-"),
        ("Remarks", sop.notes or "-"),
    ]
    table_shape = slide.shapes.add_table(len(fields), 2, PptxInches(0.3), PptxInches(0.75), PptxInches(4.4), PptxInches(0.5) * len(fields))
    table = table_shape.table
    table.columns[0].width = PptxInches(1.5)
    table.columns[1].width = PptxInches(2.9)
    for i, (label, val) in enumerate(fields):
        table.cell(i, 0).text = label
        table.cell(i, 1).text = val
        for cell in (table.cell(i, 0), table.cell(i, 1)):
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)

    proc_text, images = _split_procedure(sop.content)
    proc_box = slide.shapes.add_textbox(PptxInches(5.0), PptxInches(0.75), PptxInches(4.7), PptxInches(3.5))
    proc_tf = proc_box.text_frame
    proc_tf.word_wrap = True
    proc_tf.text = "Procedure & Method"
    proc_tf.paragraphs[0].font.bold = True
    proc_tf.paragraphs[0].font.size = Pt(12)
    for line in proc_text.split("\n"):
        if not line.strip():
            continue
        p = proc_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)

    img_top = PptxInches(4.4)
    for img_bytes in images:
        try:
            slide.shapes.add_picture(BytesIO(img_bytes), PptxInches(0.3), img_top, height=PptxInches(2.2))
            img_top += PptxInches(2.3)
        except Exception:
            pass


def generate_pptx(db: Session, sop_ids: list[int], variant: str = "절차서", cover_info: dict | None = None) -> bytes:
    sops = _get_sops_ordered(db, sop_ids)
    prs = Presentation()
    if variant == "계획서":
        _add_cover_slide_pptx(prs, cover_info)
    for sop in sops:
        _add_sop_slide_pptx(prs, sop)
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
